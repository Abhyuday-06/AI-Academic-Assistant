from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from fastapi.responses import JSONResponse, Response
from typing import Optional, List
import asyncio
import structlog
import re
import os
import tempfile
import numpy as np
from app.models import (
    NotesParseRequest, NotesParseResponse, SummarizeRequest, SummaryResponse,
    QuestionGenerationRequest, QuestionGenerationResponse, TestType, GeneratedQuestionPaper
)
from app.agents.academic_agent import academic_agent
from app.agents.question_generator import question_generator
from app.utils.config import settings
from app.utils.pdf_exporter import pdf_exporter

# Conditional imports for heavy libraries
try:
    import cv2
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

logger = structlog.get_logger()
router = APIRouter(prefix="", tags=["Academic Assistant"])


def _extract_text_from_powerpoint(file_path: str) -> str:
    """Extract text from PowerPoint files (.pptx)."""
    try:
        from pptx import Presentation
        
        text_content = ""
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)
        
        return text_content.strip()
    except ImportError:
        raise Exception("`python-pptx` is required for PowerPoint support. Please install it.")
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")


def _deskew_image(image: np.ndarray) -> np.ndarray:
    """Corrects skew in an image using OpenCV."""
    if not OPENCV_AVAILABLE or not settings.enable_advanced_ocr_preprocessing:
        return image
        
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        gray = cv2.bitwise_not(gray)
        thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
        coords = np.column_stack(np.where(thresh > 0))
        angle = cv2.minAreaRect(coords)[-1]
        
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
            
        if abs(angle) > 0.5:
            (h, w) = image.shape[:2]
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, angle, 1.0)
            rotated = cv2.warpAffine(image, M, (w, h), 
                                     flags=cv2.INTER_CUBIC, 
                                     borderMode=cv2.BORDER_REPLICATE)
            return rotated
            
        return image
    except Exception as e:
        logger.warning("Deskewing failed, returning original image.", error=str(e))
        return image


def _preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    """Applies preprocessing steps to improve OCR accuracy."""
    if not PIL_AVAILABLE:
        raise ImportError("`Pillow` is required for image processing.")

    try:
        img_np = np.array(image.convert('RGB'))

        if OPENCV_AVAILABLE and settings.enable_advanced_ocr_preprocessing:
            img_np = _deskew_image(img_np)
            gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
            denoised = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)
            binary = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                           cv2.THRESH_BINARY, 11, 2)
            return Image.fromarray(binary)
        else:
            image = image.convert('L')
            image = ImageOps.autocontrast(image)
            image = image.filter(ImageFilter.MedianFilter(size=3))
            return image

    except Exception as e:
        logger.warning("Image preprocessing failed, using original image.", error=str(e))
        return image.convert('L')


def _is_plausible_text(text: str) -> bool:
    """Check if extracted text is plausible."""
    if not text:
        return False
    
    stripped_text = text.strip()
    
    if len(stripped_text) < 15:
        return False
    
    unique_alpha = len(set(c for c in stripped_text.lower() if c.isalpha()))
    if unique_alpha < 5:
        return False
        
    words = re.findall(r'[a-zA-Z]{2,}', stripped_text)
    if len(words) < 3:
        return False
        
    return True


async def _correct_ocr_text_with_llm(text: str) -> str:
    """Uses an LLM to correct common OCR errors."""
    if not settings.enable_ocr_correction or not _is_plausible_text(text):
        return text

    logger.info("Attempting to correct OCR output with LLM...")
    try:
        correction_request = NotesParseRequest(
            content=text,
            extract_keywords=False,
            extract_concepts=False,
            extract_questions=False
        )
        
        corrected_result = await academic_agent.parse_only(correction_request)
        
        if corrected_result and corrected_result.parsed_content and len(corrected_result.parsed_content) > 0.5 * len(text):
            logger.info("LLM correction successful.")
            return corrected_result.parsed_content
        else:
            logger.warning("LLM correction returned invalid content. Returning original text.")
            return text
    except Exception as e:
        logger.error("LLM-based OCR correction failed.", error=str(e))
        return text


async def _extract_text_with_ocr(pdf_content: bytes) -> str:
    """Core OCR pipeline for PDFs."""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError as e:
        raise Exception(f"OCR library not found: {e}. Please install `pdf2image` and `pytesseract`.")

    # Configure paths for Windows
    if os.name == 'nt':
        tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        poppler_path = r"C:\poppler-24.02.0\Library\bin"
        if not os.path.exists(poppler_path):
             poppler_path = None
    else:
        poppler_path = None

    images = convert_from_bytes(
        pdf_content,
        dpi=settings.ocr_dpi,
        thread_count=settings.ocr_parallel_workers,
        poppler_path=poppler_path
    )

    if not images:
        raise Exception("Could not convert PDF to images.")

    full_text = ""
    for i, image in enumerate(images):
        logger.info(f"Processing page {i+1}/{len(images)} with OCR...")
        
        processed_image = _preprocess_image_for_ocr(image)
        config = r'--oem 3 --psm 6'
        page_text = pytesseract.image_to_string(processed_image, config=config, lang='eng')
        
        if _is_plausible_text(page_text):
            full_text += f"\n\n--- Page {i+1} ---\n{page_text}"
        else:
            logger.warning(f"Page {i+1} produced implausible text, skipping.")

    if not full_text:
        raise Exception("OCR could not extract any plausible text from the document.")
        
    corrected_text = await _correct_ocr_text_with_llm(full_text)
    return corrected_text.strip()


async def _extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from various file formats."""
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # Try direct text extraction first
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
            if _is_plausible_text(text):
                logger.info("Successfully extracted text directly from PDF.")
                return text.strip()
            else:
                logger.info("Direct text extraction yielded implausible text. Falling back to OCR.")
        except Exception as e:
            logger.warning("Direct PDF text extraction failed, falling back to OCR.", error=str(e))

        # Fallback to OCR
        logger.info("Starting OCR pipeline for PDF.")
        with open(file_path, 'rb') as f:
            pdf_content = f.read()
        return await _extract_text_with_ocr(pdf_content)

    elif file_ext == 'docx':
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            raise Exception("`python-docx` is required for DOCX support.")
            
    elif file_ext == 'pptx':
        return _extract_text_from_powerpoint(file_path)

    elif file_ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
            
    else:
        raise Exception(f"Unsupported file format: {file_ext}")


from pydantic import BaseModel
from typing import Dict, Any

# --- Request Models ---

class PDFExportRequest(BaseModel):
    export_type: str  # "summary" or "question-paper"
    data: Dict[str, Any]
    filename: Optional[str] = None


# --- Main Endpoints ---

@router.post("/generate-summary", response_model=NotesParseResponse)
async def generate_summary(
    file: UploadFile = File(None, description="Academic file (PDF, DOCX, PPTX, TXT)"),
    text: str = Form(None, description="Raw text content"),
    extract_keywords: bool = Form(True),
    extract_concepts: bool = Form(True),
    extract_questions: bool = Form(False)
):
    """
    **Main Endpoint for Note Parsing & Summarization**
    
    Accepts either a file upload OR raw text.
    Performs:
    - Text extraction (OCR for scanned PDFs)
    - Intelligent summarization
    - Keyword & Concept extraction
    """
    try:
        content = ""
        filename = "content"

        # 1. Handle File Upload
        if file:
            if not file.filename:
                raise HTTPException(status_code=400, detail="No file provided.")
            
            filename = file.filename
            file_ext = filename.lower().split('.')[-1]
            allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
            
            if file_ext not in allowed_extensions:
                raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}")

            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
                file_content = await file.read()
                if len(file_content) > settings.MAX_FILE_SIZE:
                    raise HTTPException(status_code=413, detail=f"File too large. Max size: {settings.max_file_size_mb}MB")
                temp_file.write(file_content)
                temp_path = temp_file.name

            try:
                content = await _extract_text_from_file(temp_path, filename)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

        # 2. Handle Raw Text
        elif text:
            content = text.strip()
        
        else:
            raise HTTPException(status_code=400, detail="Either 'file' or 'text' must be provided.")

        # 3. Validate Content
        if not content or len(content) < 10:
            raise HTTPException(status_code=422, detail="Could not extract meaningful content.")

        if len(content) > settings.max_content_length:
            content = content[:settings.max_content_length]
            logger.warning(f"Content truncated to {settings.max_content_length} chars.")

        # 4. Process with AI Agent
        request = NotesParseRequest(
            content=content,
            extract_keywords=extract_keywords,
            extract_concepts=extract_concepts,
            extract_questions=extract_questions
        )
        
        result = await academic_agent.parse_and_summarize(request)
        logger.info(f"Successfully processed summary for {filename}")
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Summary generation failed: {e}")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")


@router.get("/question-generation-options")
async def get_question_generation_options():
    """Get available options for the question generation form (Test Types, Modules)."""
    return {
        "test_types": [
            {"value": "CAT-1", "label": "CAT-1 (5 Questions, 50 Marks)"},
            {"value": "CAT-2", "label": "CAT-2 (5 Questions, 50 Marks)"},
            {"value": "FAT", "label": "FAT (10 Questions, 100 Marks)"}
        ],
        "modules": [
            {"value": f"Module {i}", "label": f"Module {i}"} 
            for i in range(1, 11)
        ],
        "default_modules": {
            "CAT-1": ["Module 1", "Module 2", "Module 3"],
            "CAT-2": ["Module 1", "Module 2", "Module 3"], 
            "FAT": [f"Module {i}" for i in range(1, 11)]
        }
    }


@router.post("/generate-question-paper", response_model=QuestionGenerationResponse)
async def generate_question_paper(
    file: UploadFile = File(None, description="Syllabus file (PDF, DOCX, PPTX, TXT)"),
    syllabus_text: str = Form(None, description="Raw syllabus text"),
    test_type: str = Form(..., description="CAT-1, CAT-2, or FAT"),
    modules: str = Form(..., description="Comma-separated modules (e.g., 'Module 1,Module 2')")
):
    """
    **Main Endpoint for Question Paper Generation**
    
    Accepts syllabus (File or Text) and generates a structured exam paper.
    Strictly follows the selected modules.
    """
    try:
        syllabus_content = ""

        # 1. Handle File
        if file:
            if not file.filename:
                raise HTTPException(status_code=400, detail="No file provided.")
            
            file_ext = file.filename.lower().split('.')[-1]
            allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
            if file_ext not in allowed_extensions:
                raise HTTPException(status_code=400, detail="Unsupported file type.")

            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
                file_content = await file.read()
                temp_file.write(file_content)
                temp_path = temp_file.name

            try:
                syllabus_content = await _extract_text_from_file(temp_path, file.filename)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

        # 2. Handle Text
        elif syllabus_text:
            syllabus_content = syllabus_text.strip()
        
        else:
            raise HTTPException(status_code=400, detail="Either 'file' or 'syllabus_text' must be provided.")

        # 3. Validate Content
        if not syllabus_content or len(syllabus_content) < 10:
            raise HTTPException(status_code=422, detail="Invalid syllabus content.")

        # 4. Validate Parameters
        if test_type not in ["CAT-1", "CAT-2", "FAT"]:
            raise HTTPException(status_code=400, detail="Invalid test type.")
        
        modules_list = [m.strip() for m in modules.split(',') if m.strip()]
        if not modules_list:
            raise HTTPException(status_code=400, detail="At least one module must be selected.")

        # 5. Generate Questions
        request = QuestionGenerationRequest(
            syllabus_text=syllabus_content,
            test_type=TestType(test_type),
            modules=modules_list
        )
        
        result = await question_generator.generate_question_paper(request)
        
        if not result.success:
            raise HTTPException(status_code=500, detail=result.message)
            
        logger.info(f"Successfully generated {test_type} paper.")
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Question generation failed: {e}")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")


@router.post("/export-pdf")
async def export_pdf(request: PDFExportRequest):
    """
    **Universal PDF Export Endpoint**
    
    Takes the JSON result from /generate-summary or /generate-question-paper
    and converts it into a professional PDF.
    
    - export_type: "summary" or "question-paper"
    - data: The JSON response from the previous step
    """
    try:
        pdf_content = None
        filename = request.filename or "export"

        if request.export_type == "summary":
            # Reconstruct model from JSON data
            try:
                # Handle case where data might be wrapped or raw
                if "parsed_content" in request.data:
                    model = NotesParseResponse(**request.data)
                    pdf_content = pdf_exporter.export_parse_results(model, filename)
                    filename = pdf_exporter.generate_filename(filename, "summary")
                else:
                    # Fallback for simple summary response if needed
                    model = SummaryResponse(**request.data)
                    pdf_content = pdf_exporter.export_summary_results(model, filename)
                    filename = pdf_exporter.generate_filename(filename, "summary")
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Invalid summary data format: {str(e)}")

        elif request.export_type == "question-paper":
            try:
                # We need the 'question_paper' part of the response
                if "question_paper" in request.data:
                    qp_data = request.data["question_paper"]
                else:
                    qp_data = request.data # Assume it's the paper object itself if not wrapped
                
                model = GeneratedQuestionPaper(**qp_data)
                pdf_content = pdf_exporter.export_question_paper(model, filename)
                filename = pdf_exporter.generate_filename(filename, "question_paper")
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Invalid question paper data format: {str(e)}")

        else:
            raise HTTPException(status_code=400, detail="Invalid export_type. Must be 'summary' or 'question-paper'.")

        if not pdf_content:
            raise HTTPException(status_code=500, detail="Failed to generate PDF content.")

        return Response(
            content=pdf_content,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Content-Length": str(len(pdf_content))
            }
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PDF export failed: {e}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")


@router.get("/activity")
async def get_recent_activity():
    """
    Get recent user activity (Mock Data).
    """
    return [
        {
            "id": 1,
            "type": "summary",
            "title": "Computer Networks Module 1",
            "date": "2024-01-04T10:30:00",
            "status": "completed"
        },
        {
            "id": 2,
            "type": "question_paper",
            "title": "CAT-1 Practice Paper",
            "date": "2024-01-03T14:15:00",
            "status": "completed"
        },
        {
            "id": 3,
            "type": "summary",
            "title": "Operating Systems Notes",
            "date": "2024-01-02T09:00:00",
            "status": "completed"
        }
    ]

